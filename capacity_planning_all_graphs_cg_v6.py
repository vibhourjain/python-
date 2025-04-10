import streamlit as st
import duckdb
import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

# Streamlit UI Configuration
st.set_page_config(layout="wide")
st.title("Capacity Planning Visualizer")


# PowerPoint initialization
@st.cache_data
def create_ppt_template():
    prs = Presentation()
    return prs


ppt_prs = create_ppt_template()

# Main application
reference_date = st.date_input("Select Reference Date")

if st.button("Generate Graphs"):
    con = duckdb.connect("capacity_panning.duckdb")

    # Get metadata
    metadata_df = con.execute("""
        SELECT DISTINCT LOWER(target_table) AS table_name, application
        FROM capacity_planning.query_metadata
    """).fetchdf()

    for _, row in metadata_df.iterrows():
        table_name = row["table_name"]
        app_name = row["application"]

        try:
            # Query last 13 months data
            query = f"""
                SELECT *, 
                       strftime(Period, '%Y-%m') AS month_year
                FROM {table_name}
                WHERE Period BETWEEN DATE '{reference_date}' - INTERVAL 13 MONTHS AND DATE '{reference_date}'
            """
            df = con.execute(query).fetchdf()

            if df.empty:
                st.warning(f"No data for {app_name} in past 13 months.")
                continue

            df["Period"] = pd.to_datetime(df["Period"])
            x_dates = df['Period'].drop_duplicates().sort_values()

            # Create PowerPoint slide
            slide = ppt_prs.slides.add_slide(ppt_prs.slide_layouts[6])
            slide.shapes.title.text = f"{app_name} Analysis - {reference_date}"

            # --- LINE CHART ---
            fig_line, ax_line = plt.subplots(figsize=(8, 4))
            sns.lineplot(data=df, x="Period", y="Volume",
                         hue="Application_Type", ax=ax_line, palette='colorblind')

            # X-axis tick management
            n_ticks = min(30, len(x_dates))
            tick_indices = []
            if len(x_dates) > n_ticks:
                step = len(x_dates) // n_ticks
                tick_indices = list(range(0, len(x_dates), step))[:n_ticks]
                if (len(x_dates) - 1) not in tick_indices:
                    tick_indices.append(len(x_dates) - 1)
            else:
                tick_indices = range(len(x_dates))

            xticks = x_dates.iloc[tick_indices]
            xtick_labels = [d.strftime('%Y-%m-%d') for d in xticks]

            ax_line.set_xticks(xticks)
            ax_line.set_xticklabels(xtick_labels, rotation=-45, ha='right')

            # Styling
            ax_line.spines[['top', 'right', 'left', 'bottom']].set_visible(False)
            ax_line.grid(True, axis='y', linestyle='--', alpha=0.5)
            ax_line.set_title(f"{app_name} Daily Trend")
            ax_line.legend().set_visible(False)

            # Add to PowerPoint (bottom right)
            img_line = BytesIO()
            fig_line.savefig(img_line, format='png', dpi=300, bbox_inches='tight')
            slide.shapes.add_picture(img_line, Inches(5), Inches(4), width=Inches(5))
            plt.close(fig_line)

            # --- BAR CHART ---
            fig_bar, ax_bar = plt.subplots(figsize=(8, 4))
            monthly_df = df.groupby(['month_year', 'Application_Type'], as_index=False).agg(Volume=('Volume', 'sum'))

            # Dynamic bar width
            n_categories = len(monthly_df['Application_Type'].unique())
            bar_width = 0.8 / n_categories if n_categories > 1 else 0.8

            sns.barplot(
                data=monthly_df,
                x="month_year",
                y="Volume",
                hue="Application_Type",
                ax=ax_bar,
                width=bar_width,
                palette='colorblind'
            )

            # Annotations and styling
            for container in ax_bar.containers:
                ax_bar.bar_label(container, fmt="%.0f", label_type='edge', fontsize=8)

            ax_bar.spines[['top', 'right', 'left', 'bottom']].set_visible(False)
            ax_bar.grid(True, axis='y', linestyle='--', alpha=0.5)
            ax_bar.set_title(f"{app_name} Monthly Summary")
            ax_bar.set_xticklabels(ax_bar.get_xticklabels(), rotation=-45, ha='right')
            ax_bar.legend().set_visible(False)

            # Add to PowerPoint (top left)
            img_bar = BytesIO()
            fig_bar.savefig(img_bar, format='png', dpi=300, bbox_inches='tight')
            slide.shapes.add_picture(img_bar, Inches(0.5), Inches(1.5), width=Inches(5))
            plt.close(fig_bar)

            # Streamlit display
            st.subheader(app_name)
            col1, col2 = st.columns(2)
            with col1:
                st.pyplot(fig_bar)
            with col2:
                st.pyplot(fig_line)
            st.markdown("---")

        except Exception as e:
            st.error(f"Error processing {app_name}: {str(e)}")
            continue

    # Final PowerPoint download
    ppt_buffer = BytesIO()
    ppt_prs.save(ppt_buffer)
    ppt_buffer.seek(0)

    st.download_button(
        label="📥 Download PowerPoint Report",
        data=ppt_buffer,
        file_name=f"capacity_report_{reference_date}.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

    con.close()