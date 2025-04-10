import streamlit as st
import duckdb
import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

WORK_DIR = r"C:\Users\vibho\ado-0001\python-"
DUCK_DB_PATH = r"C:\Users\vibho\ado-0001\python-\cap.duckdb"
METADATA_TABLE = "capacity_planning.query_metadata"
APPLICATION_TABLE = "capacity_planning.application"

# Streamlit UI
st.title("Capacity Planning Visualizer")

reference_date = st.date_input("Select Reference Date")
try:
    if st.button("Generate Graphs"):
        try:
            con = duckdb.connect(DUCK_DB_PATH)

            # Get metadata
            metadata_df = con.execute("""
                        SELECT DISTINCT LOWER(target_table) AS table_name, application
                        FROM capacity_planning.query_metadata where 1=1 and 2=2
                    """).fetchdf()

            for _, row in metadata_df.iterrows():
                table_name = row["table_name"]
                app_name = row["application"]

                # Query last 13 months data
                query = f"""
                            SELECT *, 
                                   strftime(Period, '%Y-%m') AS month_year
                            FROM {table_name}
                            WHERE Period BETWEEN DATE '{reference_date}' - INTERVAL 13 MONTHS AND DATE '{reference_date}'
                        """
                try:
                    df = con.execute(query).fetchdf()
                except Exception as e:
                    st.error(f"Error fetching data from {table_name}: {e}")
                    continue

                if df.empty:
                    st.warning(f"No data for {app_name} in past 13 months.")
                    continue

                df["Period"] = pd.to_datetime(df["Period"])
                x_dates = df['Period'].drop_duplicates().sort_values()

                # === LINE CHART ===
                fig, ax = plt.subplots(figsize=(14, 4))
                sns.lineplot(data=df, x="Period", y="Volume", hue="Application_Type", ax=ax, palette='muted')

                # X-axis tick management
                n_ticks = min(30, len(x_dates))
                tick_indices = []
                if len(x_dates) > n_ticks:
                    step = len(x_dates) // n_ticks
                    tick_indices = list(range(0, len(x_dates), step))[:n_ticks]
                    if (len(x_dates) - 1) not in tick_indices:
                        tick_indices.append(len(x_dates) - 1)
                else:
                    tick_indices = range(len(x_dates))

                xticks = x_dates.iloc[tick_indices]
                xtick_labels = [d.strftime('%Y/%d/%m') for d in xticks]

                ax.set_xticks(xticks)
                ax.set_xticklabels(xtick_labels, rotation=-90, ha='right')

                # Clean style
                ax.spines[['top', 'right', 'left', 'bottom']].set_visible(False)
                ax.grid(True, axis='y', linestyle='--', linewidth=0.5, color='gray', alpha=0.3)
                ax.set_axisbelow(True)
                ax.tick_params(axis='y', left=False)
                ax.tick_params(axis='x', bottom=False)
                ax.set_title(f"{app_name} Daily")
                ax.set_xlabel("Period-Daily")
                ax.set_ylabel("Volume")
                ax.legend(
                    title="Type",
                    loc='upper left',
                    bbox_to_anchor=(1.01, 1),
                    borderaxespad=0,
                    frameon=False
                )
                st.pyplot(fig)
                st.markdown("---")

                image_stream = BytesIO()
                fig.savefig(image_stream, format='png', dpi=300, bbox_inches='tight')
                ppt_template_path = 'T_pp.pptx'
                prs = Presentation(ppt_template_path)
                slide_daily = prs.slides[2]
                left = Inches(5)
                top = Inches(4)
                slide_daily.shapes.add_picture(image_stream, left, top, width=Inches(5))
                plt.close()

                # === BAR CHART ===
                monthly_df = df.groupby(['month_year', 'Application_Type'], as_index=False)['Volume'].sum()
                fig, ax = plt.subplots(figsize=(14, 10))
                sns.barplot(
                    data=monthly_df,
                    x="month_year",
                    y="Volume",
                    hue="Application_Type",
                    ax=ax,
                    width=0.5,
                    palette='colorblind'
                )

                # Extend Y limit
                max_volume = monthly_df["Volume"].max()
                ax.set_ylim(0, max_volume * 1.2)

                # Annotate bars
                for container in ax.containers:
                    ax.bar_label(container, fmt="%.0f", label_type='edge', fontsize=8)

                # Clean look
                ax.spines[['top', 'right', 'left', 'bottom']].set_visible(False)
                ax.grid(True, axis='y', linestyle='--', linewidth=0.5, color='gray', alpha=0.3)
                ax.set_axisbelow(True)
                ax.tick_params(axis='y', left=False)
                ax.tick_params(axis='x', bottom=False)

                # X ticks - rotate and format
                ax.set_xticklabels(ax.get_xticklabels(), rotation=-90, ha='right')

                ax.set_title(f"{app_name} Monthly")
                ax.set_xlabel("Period-Month-Year")
                ax.set_ylabel("Total Volume")

                # Reposition legend to outside top-right
                ax.legend(
                    title="Type",
                    loc='upper left',
                    bbox_to_anchor=(1.01, 1),
                    borderaxespad=0,
                    frameon=False
                )

                # Adjust layout for space
                plt.subplots_adjust(bottom=0.3, right=0.75)

                st.pyplot(fig)
                st.markdown("###")

                image_stream = BytesIO()
                fig.savefig(image_stream, format='png', dpi=300, bbox_inches='tight')
                slide_daily = prs.slides[2]
                left = Inches(0.5)
                top = Inches(1.5)
                slide_daily.shapes.add_picture(image_stream, left, top, width=Inches(5))
                now = datetime.now()
                q = str((int(now.strftime("%m")) % 3))
                y = now.strftime("%Y")
                pptx_out_filename = "Out_CP.pptx"
                prs.save(pptx_out_filename)
                plt.close()
        finally:
            con.close()
finally:
    pass


