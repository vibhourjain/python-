# Author Name: Vibhour Jain
# Final version with all enhancements: title, metric/peak tables, UI output
# Includes fixes for NaN values, color formatting, and slide layout.
# No Border
import streamlit as st
import duckdb
import pandas as pd
import numpy as np
import seaborn as sns
import matplotlib.pyplot as plt
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
from capacity_planning_sql_metric_v3 import insert_metrics_to_duckdb, insert_peaks_to_duckdb

WORK_DIR = r"C:\\Users\\vibho\\ado-0001\\python-"
DUCKDB_PATH = r"C:\\Users\\vibho\\ado-0001\\python-\\sybase_data.duckdb"

def gen_one_click_capacity_report():
    st.title("Capacity Planning Visualizer")

    reference_date = st.date_input("Select Reference Date")

    if st.button("Generate Graphs"):
        con = duckdb.connect(DUCKDB_PATH)
        try:
            metadata_df = con.execute("""
                SELECT DISTINCT LOWER(target_table) AS table_name, application
                FROM capacity_planning.query_metadata 
                WHERE 1=1
                AND IS_ACTIVE = 'Y'
                AND application = 'BANA-TBAR'
            """).fetchdf()

            for _, row in metadata_df.iterrows():
                table_name = row["table_name"]
                app_name = row["application"]

                query = f"""
                    SELECT *, strftime(Period, '%Y-%m') AS month_year
                    FROM {table_name}
                    WHERE Period >= (DATE '{reference_date}' - INTERVAL 13 MONTHS + INTERVAL 1 DAY)
                    AND Period < (DATE '{reference_date}' + INTERVAL 1 DAY)
                """
                try:
                    df = con.execute(query).fetchdf()
                except Exception as e:
                    st.error(f"Error fetching data from {table_name}: {e}")
                    continue

                if df.empty:
                    st.warning(f"No data for {app_name} in past 13 months.")
                    continue

                df["Period"] = pd.to_datetime(df["Period"])
                df = df.sort_values("Period").reset_index(drop=True)
                df["Period_Str"] = df["Period"].dt.strftime('%Y/%m/%d')
                df["month_year"] = df["Period"].dt.strftime("%Y-%m")
                order_hue = sorted(df['Application_Type'].unique())
                fig, ax = plt.subplots(figsize=(14, 4))
                sns.lineplot(data=df, x=df.index, y="Volume", hue="Application_Type", ax=ax, palette='muted', hue_order = order_hue)
                tick_indices = np.linspace(0, len(df) - 1, min(40, len(df)), dtype=int)
                xtick_labels = df.loc[tick_indices, "Period_Str"]

                ax.set_xticks(tick_indices)
                ax.set_xticklabels(xtick_labels, rotation=-90, ha='right')
                ax.spines[['top', 'right', 'left', 'bottom']].set_visible(False)
                ax.grid(True, axis='y', linestyle='--', linewidth=0.5, color='gray', alpha=0.5)
                ax.set_axisbelow(True)
                ax.tick_params(axis='y', left=False)
                ax.tick_params(axis='x', bottom=False)
                ax.set_title(f"{app_name} Daily")
                ax.set_xlabel("Period-Daily")
                ax.set_ylabel("Volume")
                ax.legend(title="Type", loc='upper left', bbox_to_anchor=(1.01, 1), frameon=False)
                st.pyplot(fig)
                st.markdown("---")

                image_stream = BytesIO()
                fig.savefig(image_stream, format='png', dpi=300, bbox_inches='tight')
                plt.close()

                if app_name.startswith('BANA'):
                    ppt_template_path = 'T_BANA_CapacityPerformanceReview.pptx'
                else:
                    ppt_template_path = 'T_SDI_CapacityPerformanceReview.pptx'

                prs = Presentation(ppt_template_path)
                slide_daily = prs.slides[2]

                if slide_daily.shapes.title:
                    slide_daily.shapes.title.text = f"Application Capacity Review: {app_name}"

                title_shape = slide_daily.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.5))
                tf = title_shape.text_frame
                tf.paragraphs[0].font.size = Pt(18)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

                slide_daily.shapes.add_picture(image_stream, Inches(5), Inches(4), width=Inches(5))

                monthly_df = df.groupby(['month_year', 'Application_Type'], as_index=False)['Volume'].sum()
                bar_groups = monthly_df.groupby('month_year')['Application_Type'].nunique()
                avg_bar_count = bar_groups.mean()
                bar_width = 0.5 if avg_bar_count > 1 else 0.25

                fig, ax = plt.subplots(figsize=(14, 10))
                sns.barplot(
                    data=monthly_df,
                    x="month_year",
                    y="Volume",
                    hue="Application_Type",
                    ax=ax,
                    width=bar_width,
                    palette='colorblind',
                    hue_order = order_hue
                )

                max_volume = monthly_df["Volume"].max()
                ax.set_ylim(0, max_volume * 1.2)

                for container in ax.containers:
                    ax.bar_label(container, fmt="%.0f", label_type='edge', fontsize=8)

                ax.spines[['top', 'right', 'left', 'bottom']].set_visible(False)
                ax.grid(True, axis='y', linestyle='--', linewidth=0.5, color='gray', alpha=0.5)
                ax.set_axisbelow(True)
                ax.tick_params(axis='y', left=False)
                ax.tick_params(axis='x', bottom=False)
                ax.set_xticklabels(ax.get_xticklabels(), rotation=-90, ha='right')
                ax.set_title(f"{app_name} Monthly")
                ax.set_xlabel("Period-Month-Year")
                ax.set_ylabel("Total Volume")
                ax.legend(title="Type", loc='upper left', bbox_to_anchor=(1.01, 1), frameon=False)
                plt.subplots_adjust(bottom=0.3, right=0.75)

                st.pyplot(fig)
                st.markdown("###")

                image_stream = BytesIO()
                fig.savefig(image_stream, format='png', dpi=300, bbox_inches='tight')
                slide_daily.shapes.add_picture(image_stream, Inches(0.5), Inches(1.5), width=Inches(5))
                plt.close()

                insert_metrics_to_duckdb(con, table_name, app_name, reference_date)
                insert_peaks_to_duckdb(con, table_name, app_name, reference_date)

                metric_df = con.execute(f"""
                    SELECT app_type AS "Type", metric AS "Metric",
                           baseline_13m AS "Baseline(13 Month)", current_qtr AS "Current Qtr",
                           diff_baseline_pct AS "% Diff Baseline", previous_qtr AS "Previous Qtr",
                           diff_qoq_pct AS "% Diff QoQ"
                    FROM capacity_planning.metrics_summary
                    WHERE application_name = '{app_name}' AND reference_date = '{reference_date}'
                    ORDER BY app_type, metric
                """).fetchdf()

                st.subheader(f"{app_name} Metrics Summary")
                st.dataframe(metric_df.style.format(precision=2))

                # For the metric table - create a plain table with white background, black text, and borders
                rows, cols = metric_df.shape
                metric_shape = slide_daily.shapes.add_table(rows + 1, cols, Inches(5.5), Inches(0.5), Inches(4), Inches(1.6))
                metric_table = metric_shape.table
                
                # Apply a simple table style that has borders but no background color
                # Light Style 1 - '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'
                tbl = metric_shape._element.graphic.graphicData.tbl
                style_id = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'  # Light Style 1 with borders
                tbl[0][-1].text = style_id
                
                # Set header row
                for col_idx, col_name in enumerate(metric_df.columns):
                    cell = metric_table.cell(0, col_idx)
                    cell.text = col_name
                    
                    # Set white background
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    
                    # Format text - black, bold, size 10
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = Pt(10)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0, 0, 0)

                # Fill data rows
                for row_idx, row in metric_df.iterrows():
                    for col_idx, value in enumerate(row):
                        cell = metric_table.cell(row_idx + 1, col_idx)
                        cell.text = "" if pd.isna(value) else str(value)
                        
                        # Set white background
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                        
                        # Format text - black, size 10
                        p = cell.text_frame.paragraphs[0]
                        p.font.size = Pt(10)
                        p.font.bold = False
                        p.font.color.rgb = RGBColor(0, 0, 0)

                peak_df = con.execute(f"""
                    SELECT PK_ROW FROM capacity_planning.peak_summary
                    WHERE application_name = '{app_name}' AND reference_date = '{reference_date}'
                """).fetchdf()

                if not peak_df.empty:
                    # Instead of creating a table, create a text box for peak information
                    # Adjust position to prevent overlap with metric table
                    peak_text = "\n".join(peak_df["PK_ROW"])
                    
                    # Create a text box for peak information - moved down to avoid overlap
                    peak_textbox = slide_daily.shapes.add_textbox(Inches(5.5), Inches(2.5), Inches(4), Inches(1))
                    tf = peak_textbox.text_frame
                    tf.text = peak_text
                    
                    # Format the text
                    for paragraph in tf.paragraphs:
                        paragraph.font.size = Pt(10)
                        paragraph.font.color.rgb = RGBColor(0, 0, 0)

                now = datetime.now()
                qtr_month = int(now.strftime("%m"))
                if qtr_month % 3 == 0:
                    q = (qtr_month // 3)
                else:
                    q = (qtr_month // 3) + 1

                T=""
                y = now.strftime("%Y")
                pptx_out_filename = T+ 'CapacityPerformanceReview_' + app_name + '_' + y + 'Q' + str(q) + '.pptx'
                prs.save(pptx_out_filename)

        finally:
            con.close()

        st.success("Graphs and metrics have been generated and saved to PowerPoint.")
