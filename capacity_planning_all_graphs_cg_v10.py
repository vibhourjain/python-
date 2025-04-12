import streamlit as st
import duckdb
import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt
import numpy as np
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from io import BytesIO
from capacity_panning_sql_metric_v3 import insert_metrics_to_duckdb, insert_peaks_to_duckdb

WORK_DIR = r"C:\\Users\\vibho\\ado-0001\\python-"
DUCKDB_PATH = r"C:\\Users\\vibho\\ado-0001\\python-\\cap.duckdb"

st.title("Capacity Planning Visualizer")

reference_date = st.date_input("Select Reference Date")

if st.button("Generate Graphs"):
    con = duckdb.connect(DUCKDB_PATH)
    try:
        metadata_df = con.execute("""
            SELECT DISTINCT LOWER(target_table) AS table_name, application
            FROM capacity_planning.query_metadata 
            WHERE 1=1 AND application = 'BANA-TBAR'
        """).fetchdf()

        for _, row in metadata_df.iterrows():
            table_name = row["table_name"]
            app_name = row["application"]

            query = f"""
                SELECT *, strftime(Period, '%Y-%m') AS month_year
                FROM {table_name}
                WHERE Period BETWEEN DATE '{reference_date}' - INTERVAL 13 MONTHS AND DATE '{reference_date}'
            """
            try:
                df = con.execute(query).fetchdf()
            except Exception as e:
                st.error(f"Error fetching data from {table_name}: {e}")
                continue

            if df.empty:
                st.warning(f"No data for {app_name} in past 13 months.")
                continue

            df["Period"] = pd.to_datetime(df["Period"])
            df = df.sort_values("Period").reset_index(drop=True)
            df["Period_Str"] = df["Period"].dt.strftime('%Y/%m/%d')
            df["month_year"] = df["Period"].dt.strftime("%Y-%m")

            ppt_template_path = 'T_BANA_CapacityPerformanceReview.pptx'
            prs = Presentation(ppt_template_path)
            slide_daily = prs.slides[2]

            # Bar chart
            monthly_df = df.groupby(['month_year', 'Application_Type'], as_index=False)['Volume'].sum()
            fig, ax = plt.subplots(figsize=(14, 6))
            sns.barplot(
                data=monthly_df,
                x="month_year",
                y="Volume",
                hue="Application_Type",
                ax=ax,
                palette='colorblind'
            )
            ax.set_xticklabels(ax.get_xticklabels(), rotation=-90, ha='right')
            ax.set_title(f"{app_name} Monthly")
            ax.set_xlabel("Month-Year")
            ax.set_ylabel("Volume")
            plt.tight_layout()
            bar_stream = BytesIO()
            fig.savefig(bar_stream, format='png', dpi=300, bbox_inches='tight')
            slide_daily.shapes.add_picture(bar_stream, Inches(0.5), Inches(0.5), width=Inches(5))
            plt.close()

            # Line chart
            fig, ax = plt.subplots(figsize=(14, 4))
            sns.lineplot(data=df, x=df.index, y="Volume", hue="Application_Type", ax=ax, palette='muted')
            tick_indices = np.linspace(0, len(df) - 1, min(15, len(df)), dtype=int)
            ax.set_xticks(tick_indices)
            ax.set_xticklabels(df.loc[tick_indices, "Period_Str"], rotation=-90, ha='right')
            ax.set_title(f"{app_name} Daily")
            ax.set_xlabel("Period")
            ax.set_ylabel("Volume")
            plt.tight_layout()
            line_stream = BytesIO()
            fig.savefig(line_stream, format='png', dpi=300, bbox_inches='tight')
            slide_daily.shapes.add_picture(line_stream, Inches(0.5), Inches(3.3), width=Inches(5))
            plt.close()

            # Insert new metrics and peaks
            insert_metrics_to_duckdb(con, table_name, app_name, reference_date)
            insert_peaks_to_duckdb(con, table_name, app_name, reference_date)

            metric_df = con.execute(f"""
                SELECT "Type", "Metric", "Baseline(13 Month)", "Current Qtr", "% Diff Baseline", "Previous Qtr", "% Diff QoQ"
                FROM capacity_planning.metrics_summary
                WHERE application_name = '{app_name}' AND reference_date = '{reference_date}'
                ORDER BY "Type", "Metric"
            """).fetchdf()

            peak_df = con.execute(f"""
                SELECT PK_ROW
                FROM capacity_planning.peak_summary
                WHERE application_name = '{app_name}' AND reference_date = '{reference_date}'
            """).fetchdf()

            # Metric Table (Top Right)
            rows, cols = metric_df.shape
            table_shape = slide_daily.shapes.add_table(rows + 1, cols, Inches(5.5), Inches(0.5), Inches(4), Inches(1.6)).table
            for col_idx, col_name in enumerate(metric_df.columns):
                table_shape.cell(0, col_idx).text = col_name
            for row_idx, row in metric_df.iterrows():
                for col_idx, value in enumerate(row):
                    table_shape.cell(row_idx + 1, col_idx).text = str(value)
            for row in table_shape.rows:
                for cell in row.cells:
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = Pt(9)
                    p.alignment = PP_ALIGN.LEFT

            # Peak Table (Below Metric)
            if not peak_df.empty:
                peak_table = slide_daily.shapes.add_table(len(peak_df) + 1, 1, Inches(5.5), Inches(2.3), Inches(4), Inches(1)).table
                peak_table.cell(0, 0).text = "Peak Summary"
                peak_table.cell(0, 0).text_frame.paragraphs[0].font.bold = True
                for i, val in enumerate(peak_df["PK_ROW"], start=1):
                    peak_table.cell(i, 0).text = val
                    peak_table.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(9)

            # Save final output
            now = datetime.now()
            q = str((int(now.strftime("%m")) % 3))
            y = now.strftime("%Y")
            pptx_out_filename = 'CapacityPerformanceReview_' + app_name + '_Q' + y + q + '.pptx'
            prs.save(pptx_out_filename)
    finally:
        con.close()

    st.success("Graphs and metrics have been generated and saved to PowerPoint.")
