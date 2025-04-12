import streamlit as st
import duckdb
import pandas as pd
import numpy as np
import seaborn as sns
import matplotlib.pyplot as plt
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
from capacity_panning_sql_metric_v3 import insert_metrics_to_duckdb, insert_peaks_to_duckdb

WORK_DIR = r"C:\\Users\\vibho\\ado-0001\\python-"
DUCKDB_PATH = r"C:\\Users\\vibho\\ado-0001\\python-\\cap.duckdb"

def gen_one_click_capacity_report():
    st.title("Capacity Planning Visualizer")

    reference_date = st.date_input("Select Reference Date")

    if st.button("Generate Graphs"):
        con = duckdb.connect(DUCKDB_PATH)
        try:
            metadata_df = con.execute("""
                SELECT DISTINCT LOWER(target_table) AS table_name, application
                FROM capacity_planning.query_metadata 
                WHERE 1=1 AND application = 'BANA-TBAR'
            """).fetchdf()

            for _, row in metadata_df.iterrows():
                table_name = row["table_name"]
                app_name = row["application"]

                query = f"""
                    SELECT *, strftime(Period, '%Y-%m') AS month_year
                    FROM {table_name}
                    WHERE Period BETWEEN DATE '{reference_date}' - INTERVAL 13 MONTHS AND DATE '{reference_date}'
                """
                try:
                    df = con.execute(query).fetchdf()
                except Exception as e:
                    st.error(f"Error fetching data from {table_name}: {e}")
                    continue

                if df.empty:
                    st.warning(f"No data for {app_name} in past 13 months.")
                    continue

                df["Period"] = pd.to_datetime(df["Period"])
                df = df.sort_values("Period").reset_index(drop=True)
                df["Period_Str"] = df["Period"].dt.strftime('%Y/%m/%d')
                df["month_year"] = df["Period"].dt.strftime("%Y-%m")

                # Line chart with evenly spaced ticks
                fig, ax = plt.subplots(figsize=(14, 4))
                sns.lineplot(data=df, x=df.index, y="Volume", hue="Application_Type", ax=ax, palette='muted')

                n_ticks = min(15, len(df))
                tick_indices = np.linspace(0, len(df) - 1, n_ticks, dtype=int)
                xtick_labels = df.loc[tick_indices, "Period_Str"]

                ax.set_xticks(tick_indices)
                ax.set_xticklabels(xtick_labels, rotation=-90, ha='right')

                ax.spines[['top', 'right', 'left', 'bottom']].set_visible(False)
                ax.grid(True, axis='y', linestyle='--', linewidth=0.5, color='gray', alpha=0.3)
                ax.set_axisbelow(True)
                ax.tick_params(axis='y', left=False)
                ax.tick_params(axis='x', bottom=False)
                ax.set_title(f"{app_name} Daily")
                ax.set_xlabel("Period-Daily")
                ax.set_ylabel("Volume")
                ax.legend(title="Type", loc='upper left', bbox_to_anchor=(1.01, 1), frameon=False)
                st.pyplot(fig)
                st.markdown("---")

                image_stream = BytesIO()
                fig.savefig(image_stream, format='png', dpi=300, bbox_inches='tight')
                plt.close()

                if app_name.startswith('BANA'):
                    ppt_template_path = 'T_BANA_CapacityPerformanceReview.pptx'
                else:
                    ppt_template_path = 'T_SDI_CapacityPerformanceReview.pptx'

                prs = Presentation(ppt_template_path)
                slide_daily = prs.slides[2]
                slide_daily.shapes.add_picture(image_stream, Inches(5), Inches(4), width=Inches(5))

                # Bar chart
                monthly_df = df.groupby(['month_year', 'Application_Type'], as_index=False)['Volume'].sum()
                bar_groups = monthly_df.groupby('month_year')['Application_Type'].nunique()
                avg_bar_count = bar_groups.mean()
                bar_width = 0.5 if avg_bar_count > 1 else 0.25

                fig, ax = plt.subplots(figsize=(14, 10))
                sns.barplot(
                    data=monthly_df,
                    x="month_year",
                    y="Volume",
                    hue="Application_Type",
                    ax=ax,
                    width=bar_width,
                    palette='colorblind'
                )

                max_volume = monthly_df["Volume"].max()
                ax.set_ylim(0, max_volume * 1.2)

                for container in ax.containers:
                    ax.bar_label(container, fmt="%.0f", label_type='edge', fontsize=8)

                ax.spines[['top', 'right', 'left', 'bottom']].set_visible(False)
                ax.grid(True, axis='y', linestyle='--', linewidth=0.5, color='gray', alpha=0.3)
                ax.set_axisbelow(True)
                ax.tick_params(axis='y', left=False)
                ax.tick_params(axis='x', bottom=False)

                ax.set_xticklabels(ax.get_xticklabels(), rotation=-90, ha='right')
                ax.set_title(f"{app_name} Monthly")
                ax.set_xlabel("Period-Month-Year")
                ax.set_ylabel("Total Volume")
                ax.legend(title="Type", loc='upper left', bbox_to_anchor=(1.01, 1), frameon=False)
                plt.subplots_adjust(bottom=0.3, right=0.75)

                st.pyplot(fig)
                st.markdown("###")

                image_stream = BytesIO()
                fig.savefig(image_stream, format='png', dpi=300, bbox_inches='tight')
                slide_daily.shapes.add_picture(image_stream, Inches(0.5), Inches(1.5), width=Inches(5))
                plt.close()

                # Insert new metrics and peaks
                insert_metrics_to_duckdb(con, table_name, app_name, reference_date)
                insert_peaks_to_duckdb(con, table_name, app_name, reference_date)

                metric_df = con.execute(f"""
                                    SELECT app_type AS "Type", metric AS "Metric",
                                           baseline_13m AS "Baseline(13 Month)", current_qtr AS "Current Qtr",
                                           diff_baseline_pct AS "% Diff Baseline", previous_qtr AS "Previous Qtr",
                                           diff_qoq_pct AS "% Diff QoQ"
                                    FROM capacity_planning.metrics_summary
                                    WHERE application_name = '{app_name}' AND reference_date = '{reference_date}'
                                    ORDER BY app_type, metric
                                """).fetchdf()

                rows, cols = metric_df.shape
                metric_table = slide_daily.shapes.add_table(rows + 1, cols, Inches(5.5), Inches(0.5), Inches(4),
                                                            Inches(1.6)).table
                for col_idx, col_name in enumerate(metric_df.columns):
                    cell = metric_table.cell(0, col_idx)
                    cell.text = col_name
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255,255,255)
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = Pt(10)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0, 0, 0)

                for row_idx, row in metric_df.iterrows():
                    for col_idx, value in enumerate(row):
                        cell = metric_table.cell(row_idx + 1, col_idx)
                        cell.text = str(value)
                        p = cell.text_frame.paragraphs[0]
                        p.font.size = Pt(10)
                        p.font.color.rgb = RGBColor(0, 0, 0)

                # Peak summary below metric
                peak_df = con.execute(f"""
                                    SELECT PK_ROW FROM capacity_planning.peak_summary
                                    WHERE application_name = '{app_name}' AND reference_date = '{reference_date}'
                                """).fetchdf()

                # Peak Table (Below Metric)
                if not peak_df.empty:
                    peak_table = slide_daily.shapes.add_table(len(peak_df) + 1, 1, Inches(5.5), Inches(2.3), Inches(4),
                                                              Inches(1)).table
                    peak_table.cell(0, 0).text = "Peak Summary"
                    peak_table.cell(0, 0).text_frame.paragraphs[0].font.bold = True
                    for i, val in enumerate(peak_df["PK_ROW"], start=1):
                        peak_table.cell(i, 0).text = val
                        peak_table.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(9)
                        peak_table.cell(i, 0).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

                now = datetime.now()
                q = str((int(now.strftime("%m")) % 3))
                y = now.strftime("%Y")
                pptx_out_filename = 'CapacityPerformanceReview_' + app_name + '_Q' + y + q + '.pptx'
                prs.save(pptx_out_filename)
        finally:
            con.close()

        st.success("Graphs and metrics have been generated and saved to PowerPoint.")
